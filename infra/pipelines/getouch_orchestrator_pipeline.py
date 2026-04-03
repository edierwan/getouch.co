"""
title: Getouch Multimodal Orchestrator Pipeline
author: Getouch
description: Routed execution for text, documents, image understanding, and mixed multimodal requests. Includes travel/place image cards.
version: 1.2.0
requirements: pydantic,pypdf,python-docx,pandas,openpyxl,python-pptx
"""

import base64
import csv
import hashlib
import io
import json
import math
import os
import re
import sqlite3
import threading
import time
import traceback
try:
    import imghdr
except ModuleNotFoundError:
    # imghdr removed in Python 3.13+; provide minimal fallback
    class _ImghdrShim:
        _SIGS = {b'\x89PNG': 'png', b'\xff\xd8\xff': 'jpeg', b'GIF8': 'gif', b'RIFF': 'webp', b'BM': 'bmp'}
        @staticmethod
        def what(filename, h=None):
            data = h or b''
            if filename and not data:
                with open(filename, 'rb') as f:
                    data = f.read(32)
            for sig, fmt in _ImghdrShim._SIGS.items():
                if data[:len(sig)] == sig:
                    return fmt
            return None
    imghdr = _ImghdrShim()
from typing import Any, Dict, Iterator, List, Optional, Tuple, Union
from urllib.parse import quote_plus, urlparse
from urllib import request
from urllib.error import HTTPError

from pydantic import BaseModel

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


SUPPORTED_DOC_MIME = {
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/csv",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

SUPPORTED_IMAGE_MIME = {
    "image/png",
    "image/jpeg",
    "image/jpg",
    "image/webp",
    "image/gif",
    "image/bmp",
}

MIME_BY_EXT = {
    ".pdf": "application/pdf",
    ".doc": "application/msword",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".xls": "application/vnd.ms-excel",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".ppt": "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class NormalizedFile:
    def __init__(
        self,
        file_id: str,
        name: str,
        mime: str,
        size: int,
        content: bytes,
        source_hint: str,
    ):
        self.file_id = file_id
        self.name = name
        self.mime = mime
        self.size = size
        self.content = content
        self.source_hint = source_hint


class Pipeline:
    class Valves(BaseModel):
        TEXT_MODEL: str = "qwen2.5:14b"
        DOC_MODEL: str = "qwen2.5:14b"
        VISION_MODEL: str = "qwen2.5vl:32b"
        FALLBACK_MODEL: str = "qwen2.5:14b"
        OLLAMA_BASE_URL: str = "http://ollama:11434"
        SEARXNG_BASE_URL: str = "http://searxng:8080"
        OPENWEBUI_BASE_URL: str = "http://open-webui:8080"
        EMBEDDING_MODEL: str = "all-minilm:latest"
        MAX_FILE_SIZE_MB: int = 25
        MAX_FILES_PER_TURN: int = 8
        CHUNK_SIZE: int = 1200
        CHUNK_OVERLAP: int = 150
        RETRIEVAL_TOP_K: int = 6
        WEB_IMAGE_LIMIT: int = 4
        PLACE_IMAGE_LIMIT: int = 6
        MAX_IMAGES_PER_TURN: int = 8
        ENABLE_STRUCTURED_LOGS: bool = True

    def __init__(self):
        self.type = "manifold"
        self.name = "Getouch Orchestrator"
        self.pipelines = [
            {"id": "assistant", "name": "Getouch Smart Assistant"}
        ]
        self.valves = self.Valves(
            **{k: os.getenv(k, v.default) for k, v in self.Valves.model_fields.items()}
        )
        self.db_path = "/app/pipelines/getouch_orchestrator.db"
        self.lock = threading.Lock()
        self._ensure_db()

    async def on_startup(self):
        self._log("startup", {"pipeline": self.name})

    async def on_shutdown(self):
        self._log("shutdown", {"pipeline": self.name})

    async def on_valves_updated(self):
        self._log("valves_updated", self.valves.model_dump())

    def _prepare_routed_body(self, body: dict, user: Optional[dict] = None) -> dict:
        started = time.time()
        body = body or {}
        messages = body.get("messages", [])
        user_text = self._last_user_text(messages)

        attachments = self._extract_attachments(body, messages)
        normalized_docs, parse_notes = self._normalize_document_files(attachments.get("files", []))
        image_payloads = attachments.get("images", [])

        mode = self._classify_mode(user_text, normalized_docs, image_payloads)
        if mode == "tool_required":
            target_model = self.valves.TEXT_MODEL
        elif mode == "text_only":
            target_model = self.valves.TEXT_MODEL
        elif mode == "text_with_documents":
            target_model = self.valves.DOC_MODEL
        elif mode in {"image_understanding", "mixed_multimodal"}:
            target_model = self.valves.VISION_MODEL
        else:
            target_model = self.valves.FALLBACK_MODEL

        file_refs: List[str] = []
        retrieval_context = ""
        uncertainty_notes: List[str] = []

        if normalized_docs:
            for doc in normalized_docs:
                status, note = self._ingest_document(doc)
                file_refs.append(f"{doc.name} ({status})")
                if note:
                    uncertainty_notes.append(note)

            hits = self._retrieve_context(user_text, top_k=self.valves.RETRIEVAL_TOP_K)
            if hits:
                lines = []
                for h in hits:
                    lines.append(
                        f"- [{h['source']}] {h['text'][:550].replace(chr(10), ' ')}"
                    )
                retrieval_context = "\n".join(lines)

        vision_context = ""
        if image_payloads and mode in {"image_understanding", "mixed_multimodal"}:
            vision_context = "Image attachments were provided. Prioritize visible evidence from the images in your answer."

        routed_system_prompt = self._compose_system_context(
            mode=mode,
            retrieval_context=retrieval_context,
            vision_context=vision_context,
            file_refs=file_refs,
            parse_notes=parse_notes,
            uncertainty_notes=uncertainty_notes,
        )

        body.setdefault("messages", [])
        body["messages"] = [{"role": "system", "content": routed_system_prompt}] + body["messages"]
        body["model"] = target_model

        meta = body.get("metadata") or {}
        meta["getouch_route"] = {
            "mode": mode,
            "target_model": target_model,
            "doc_count": len(normalized_docs),
            "image_count": len(image_payloads),
        }
        body["metadata"] = meta

        self._log(
            "routing_decision",
            {
                "mode": mode,
                "target_model": target_model,
                "doc_count": len(normalized_docs),
                "image_count": len(image_payloads),
                "latency_ms": int((time.time() - started) * 1000),
                "user_id": (user or {}).get("id"),
            },
        )
        return body

    async def inlet(self, body: dict, user: Optional[dict] = None) -> dict:
        # Kept for compatibility when this pipeline is attached as a filter.
        return self._prepare_routed_body(body, user)

    def pipe(
        self, user_message: str, model_id: str, messages: List[dict], body: dict
    ) -> Union[str, Iterator[str]]:
        started = time.time()
        req_body = dict(body or {})
        req_body["messages"] = messages or req_body.get("messages", [])

        routed = self._prepare_routed_body(req_body, req_body.get("user"))
        target_model = routed.get("model", self.valves.FALLBACK_MODEL)
        route_meta = (routed.get("metadata") or {}).get("getouch_route") or {}
        mode = str(route_meta.get("mode") or "fallback")

        ollama_messages = self._to_ollama_messages(routed.get("messages", []))
        if not ollama_messages:
            ollama_messages = [{"role": "user", "content": user_message or "Help the user."}]

        # ── Travel itinerary planner (enhanced flow) ────────────
        if self._is_travel_itinerary(user_message) and mode in {"text_only", "tool_required"}:
            return self._handle_travel_itinerary(
                user_message=user_message,
                target_model=target_model,
                ollama_messages=ollama_messages,
                body=body,
                started=started,
            )

        place_cards: List[Dict[str, str]] = []
        web_image_urls: List[str] = []
        if self._should_search_place_images(user_message, mode):
            place_cards = self._build_place_cards(user_message)
        elif self._should_search_web_images(user_message, mode):
            web_image_urls = self._search_web_images(user_message, self.valves.WEB_IMAGE_LIMIT)

        stream_requested = bool(body.get("stream", False))
        stream_supported = mode not in {"image_understanding", "mixed_multimodal"}
        options = (routed.get("options") or {}) if isinstance(routed.get("options"), dict) else {}

        # For document modes, add num_predict to avoid infinite generation
        if mode == "text_with_documents" and "num_predict" not in options:
            options["num_predict"] = 4096

        if stream_requested and stream_supported:
            return self._stream_ollama_chat(
                target_model=target_model,
                messages=ollama_messages,
                options=options,
                web_image_urls=web_image_urls,
                place_cards=place_cards,
            )

        try:
            timeout_sec = 600 if mode in {"image_understanding", "mixed_multimodal"} else 240
            content = self._complete_ollama_chat(
                target_model=target_model,
                messages=ollama_messages,
                options=options,
                timeout=timeout_sec,
            )
            if not isinstance(content, str) or not content.strip():
                content = "I couldn't produce a grounded answer for this request. Please try a smaller file or narrower question."

            if place_cards:
                img_md = self._format_place_cards(place_cards)
                content = f"{img_md}\n\n{content}"
            elif web_image_urls:
                img_md = self._format_image_markdown(web_image_urls)
                content = f"{img_md}\n\n{content}"

            self._log(
                "model_latency",
                {
                    "target_model": target_model,
                    "latency_ms": int((time.time() - started) * 1000),
                },
            )
            return content
        except Exception as exc:
            self._log(
                "failure",
                {
                    "reason": str(exc),
                    "target_model": target_model,
                },
            )
            return (
                "I hit a routing/runtime error while processing this request. "
                "Please retry. If this continues, reduce file size or attachment count."
            )

    async def outlet(self, body: dict, user: Optional[dict] = None) -> dict:
        return body

    def _ensure_db(self):
        with self.lock:
            con = sqlite3.connect(self.db_path)
            cur = con.cursor()
            cur.execute(
                """
                CREATE TABLE IF NOT EXISTS files (
                    file_hash TEXT PRIMARY KEY,
                    file_id TEXT,
                    name TEXT,
                    mime TEXT,
                    size INTEGER,
                    extracted_text TEXT,
                    extraction_status TEXT,
                    confidence REAL,
                    parser_notes TEXT,
                    created_at INTEGER,
                    updated_at INTEGER
                )
                """
            )
            cur.execute(
                """
                CREATE TABLE IF NOT EXISTS chunks (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    file_hash TEXT,
                    source_ref TEXT,
                    chunk_index INTEGER,
                    chunk_text TEXT,
                    embedding_json TEXT,
                    confidence REAL,
                    created_at INTEGER
                )
                """
            )
            con.commit()
            con.close()

    def _extract_attachments(self, body: Dict[str, Any], messages: List[Dict[str, Any]]) -> Dict[str, Any]:
        files: List[Dict[str, Any]] = []
        images: List[str] = []

        files.extend(self._coerce_list(body.get("files")))
        files.extend(self._coerce_list(body.get("attachments")))

        for message in messages:
            if message.get("role") != "user":
                continue
            files.extend(self._coerce_list(message.get("files")))
            files.extend(self._coerce_list(message.get("attachments")))
            images.extend(self._coerce_list(message.get("images")))

            content = message.get("content")
            if isinstance(content, list):
                for part in content:
                    if not isinstance(part, dict):
                        continue
                    ptype = str(part.get("type", "")).lower()
                    if ptype in {"image_url", "input_image", "image"}:
                        img_url = (
                            part.get("image_url", {}).get("url")
                            if isinstance(part.get("image_url"), dict)
                            else part.get("image")
                        )
                        if img_url:
                            images.append(img_url)
                    if ptype in {"file", "input_file"}:
                        files.append(part)
        return {"files": files, "images": images}

    def _normalize_document_files(
        self, raw_files: List[Dict[str, Any]]
    ) -> Tuple[List[NormalizedFile], List[str]]:
        notes: List[str] = []
        normalized: List[NormalizedFile] = []
        max_bytes = self.valves.MAX_FILE_SIZE_MB * 1024 * 1024

        for idx, raw in enumerate(raw_files[: self.valves.MAX_FILES_PER_TURN]):
            try:
                item = raw if isinstance(raw, dict) else {"name": f"file-{idx}", "content": str(raw)}
                name = str(item.get("name") or item.get("filename") or f"file-{idx}")
                ext = os.path.splitext(name)[1].lower()
                mime = str(item.get("mime_type") or item.get("type") or MIME_BY_EXT.get(ext, ""))

                content = self._extract_file_bytes(item)
                size = len(content)
                if size == 0:
                    notes.append(f"Skipped {name}: empty or inaccessible content.")
                    continue
                if size > max_bytes:
                    notes.append(
                        f"Skipped {name}: exceeds max size {self.valves.MAX_FILE_SIZE_MB}MB."
                    )
                    continue
                if mime and mime not in SUPPORTED_DOC_MIME:
                    continue
                if ext and MIME_BY_EXT.get(ext) and MIME_BY_EXT[ext] not in SUPPORTED_DOC_MIME:
                    continue

                file_id = str(item.get("id") or item.get("file_id") or hashlib.sha1((name + str(size)).encode()).hexdigest())
                normalized.append(
                    NormalizedFile(
                        file_id=file_id,
                        name=name,
                        mime=mime or MIME_BY_EXT.get(ext, "application/octet-stream"),
                        size=size,
                        content=content,
                        source_hint=str(item.get("source") or "chat_attachment"),
                    )
                )
            except Exception as exc:
                notes.append(f"Failed to normalize attachment {idx}: {exc}")
        return normalized, notes

    def _extract_file_bytes(self, item: Dict[str, Any]) -> bytes:
        if isinstance(item.get("content"), str):
            c = item["content"]
            if c.startswith("data:") and ";base64," in c:
                return base64.b64decode(c.split(";base64,", 1)[1])
            return c.encode("utf-8", errors="ignore")

        if isinstance(item.get("data"), str):
            d = item["data"]
            if d.startswith("data:") and ";base64," in d:
                d = d.split(";base64,", 1)[1]
            try:
                return base64.b64decode(d)
            except Exception:
                return d.encode("utf-8", errors="ignore")

        if isinstance(item.get("bytes"), (bytes, bytearray)):
            return bytes(item.get("bytes"))

        if isinstance(item.get("path"), str) and os.path.exists(item["path"]):
            with open(item["path"], "rb") as f:
                return f.read()

        if isinstance(item.get("url"), str):
            try:
                req = request.Request(item["url"], headers={"User-Agent": "Mozilla/5.0"})
                return request.urlopen(req, timeout=20).read()
            except Exception:
                return b""
        return b""

    def _classify_mode(
        self, user_text: str, docs: List[NormalizedFile], images: List[str]
    ) -> str:
        has_docs = len(docs) > 0
        has_images = len(images) > 0

        tool_kw = re.search(
            r"\b(search|latest|current|real[- ]?time|browse|web|price|weather|stock)\b",
            user_text.lower(),
        )

        if has_docs and has_images:
            return "mixed_multimodal"
        if has_images:
            return "image_understanding"
        if has_docs:
            return "text_with_documents"
        if tool_kw:
            return "tool_required"
        if not user_text.strip():
            return "fallback"
        return "text_only"

    def _ingest_document(self, doc: NormalizedFile) -> Tuple[str, str]:
        file_hash = hashlib.sha1(doc.content).hexdigest()
        now = int(time.time())

        existing = self._db_fetchone(
            "SELECT extraction_status, parser_notes FROM files WHERE file_hash = ?",
            (file_hash,),
        )
        if existing and existing[0] == "ready":
            return "indexed", ""

        text, refs, confidence, parser_note = self._parse_document(doc)

        if not text.strip():
            status = "low_confidence"
            parser_note = (parser_note + " | No text extracted; vision/OCR fallback recommended.").strip(" |")
            confidence = min(confidence, 0.2)
        else:
            status = "ready"

        self._db_execute(
            """
            INSERT OR REPLACE INTO files
            (file_hash, file_id, name, mime, size, extracted_text, extraction_status, confidence, parser_notes, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, COALESCE((SELECT created_at FROM files WHERE file_hash = ?), ?), ?)
            """,
            (
                file_hash,
                doc.file_id,
                doc.name,
                doc.mime,
                doc.size,
                text,
                status,
                confidence,
                parser_note,
                file_hash,
                now,
                now,
            ),
        )

        self._db_execute("DELETE FROM chunks WHERE file_hash = ?", (file_hash,))

        if text.strip():
            chunks = self._chunk_text(text, self.valves.CHUNK_SIZE, self.valves.CHUNK_OVERLAP)
            for idx, chunk in enumerate(chunks):
                source_ref = refs.get(idx, refs.get(-1, doc.name))
                emb = self._embed_text(chunk)
                self._db_execute(
                    """
                    INSERT INTO chunks (file_hash, source_ref, chunk_index, chunk_text, embedding_json, confidence, created_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        file_hash,
                        source_ref,
                        idx,
                        chunk,
                        json.dumps(emb) if emb else "[]",
                        confidence,
                        now,
                    ),
                )
        return status, parser_note

    def _parse_document(
        self, doc: NormalizedFile
    ) -> Tuple[str, Dict[int, str], float, str]:
        name_lower = doc.name.lower()
        refs: Dict[int, str] = {-1: doc.name}
        note = ""
        confidence = 0.85

        try:
            if name_lower.endswith(".pdf") or doc.mime == "application/pdf":
                if PdfReader is None:
                    return "", refs, 0.1, "pypdf unavailable"
                reader = PdfReader(io.BytesIO(doc.content))
                parts = []
                chunk_idx = 0
                for i, page in enumerate(reader.pages, start=1):
                    page_text = (page.extract_text() or "").strip()
                    if page_text:
                        parts.append(page_text)
                        refs[chunk_idx] = f"{doc.name} p.{i}"
                        chunk_idx += 1
                if not parts:
                    confidence = 0.25
                    note = "PDF appears scanned/image-only"
                return "\n\n".join(parts), refs, confidence, note

            if name_lower.endswith(".docx") or "wordprocessingml" in doc.mime:
                if docx is None:
                    return "", refs, 0.1, "python-docx unavailable"
                d = docx.Document(io.BytesIO(doc.content))
                lines = []
                for p in d.paragraphs:
                    t = p.text.strip()
                    if t:
                        lines.append(t)
                for ti, table in enumerate(d.tables, start=1):
                    rows = []
                    for row in table.rows[:30]:
                        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                    if rows:
                        lines.append(f"[Table {ti}]\n" + "\n".join(rows))
                return "\n".join(lines), refs, confidence, note

            if name_lower.endswith(".txt") or name_lower.endswith(".md") or doc.mime in {"text/plain", "text/markdown"}:
                return doc.content.decode("utf-8", errors="ignore"), refs, confidence, note

            if name_lower.endswith(".csv") or "csv" in doc.mime:
                decoded = doc.content.decode("utf-8", errors="ignore")
                rdr = csv.reader(io.StringIO(decoded))
                rows = [row for _, row in zip(range(120), rdr)]
                if not rows:
                    return "", refs, 0.2, "CSV has no readable rows"
                headers = rows[0]
                preview = rows[1:21]
                summary = [
                    f"CSV file: {doc.name}",
                    f"Headers: {', '.join(headers[:40])}",
                    f"Preview rows: {len(preview)}",
                ]
                formatted = [", ".join(r[:40]) for r in preview]
                return "\n".join(summary + ["Rows:"] + formatted), refs, confidence, note

            if name_lower.endswith(".xlsx") or "spreadsheetml" in doc.mime or name_lower.endswith(".xls"):
                if openpyxl is None:
                    return "", refs, 0.1, "openpyxl unavailable"
                wb = openpyxl.load_workbook(io.BytesIO(doc.content), data_only=True, read_only=True)
                lines = [f"Workbook: {doc.name}"]
                for ws in wb.worksheets[:10]:
                    lines.append(f"Sheet: {ws.title}")
                    sample_rows = []
                    for ridx, row in enumerate(ws.iter_rows(min_row=1, max_row=20, values_only=True), start=1):
                        vals = ["" if v is None else str(v) for v in row[:20]]
                        sample_rows.append(f"R{ridx}: " + " | ".join(vals))
                    lines.extend(sample_rows)
                return "\n".join(lines), refs, confidence, note

            if name_lower.endswith(".pptx") or "presentationml" in doc.mime or name_lower.endswith(".ppt"):
                if Presentation is None:
                    return "", refs, 0.1, "python-pptx unavailable"
                prs = Presentation(io.BytesIO(doc.content))
                lines = [f"Presentation: {doc.name}"]
                for idx, slide in enumerate(prs.slides, start=1):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            texts.append(shape.text.strip())
                    if hasattr(slide, "notes_slide") and slide.notes_slide and slide.notes_slide.notes_text_frame:
                        nt = slide.notes_slide.notes_text_frame.text.strip()
                        if nt:
                            texts.append(f"Notes: {nt}")
                    lines.append(f"Slide {idx}: " + " | ".join([t for t in texts if t]))
                    refs[idx - 1] = f"{doc.name} slide {idx}"
                return "\n".join(lines), refs, confidence, note

            return "", refs, 0.15, "Unsupported document format"
        except Exception as exc:
            return "", refs, 0.1, f"Parser error: {exc}"

    def _chunk_text(self, text: str, size: int, overlap: int) -> List[str]:
        clean = re.sub(r"\s+", " ", text).strip()
        if not clean:
            return []
        chunks = []
        start = 0
        while start < len(clean):
            end = min(len(clean), start + size)
            chunk = clean[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(clean):
                break
            start = max(0, end - overlap)
        return chunks

    def _embed_text(self, text: str) -> List[float]:
        payload = {
            "model": self.valves.EMBEDDING_MODEL,
            "prompt": text[:6000],
        }
        try:
            data = self._http_json(
                f"{self.valves.OLLAMA_BASE_URL}/api/embeddings", payload, timeout=20
            )
            emb = data.get("embedding") or []
            return emb if isinstance(emb, list) else []
        except Exception:
            return []

    def _retrieve_context(self, query: str, top_k: int = 6) -> List[Dict[str, Any]]:
        if not query.strip():
            return []
        q_emb = self._embed_text(query)
        rows = self._db_fetchall(
            "SELECT source_ref, chunk_text, embedding_json, confidence FROM chunks ORDER BY created_at DESC LIMIT 500",
            (),
        )

        scored = []
        for source_ref, chunk_text, emb_json, conf in rows:
            score = 0.0
            if q_emb:
                try:
                    c_emb = json.loads(emb_json or "[]")
                    score = self._cosine(q_emb, c_emb)
                except Exception:
                    score = 0.0
            if score <= 0.0:
                score = self._keyword_overlap(query, chunk_text)
            score = score * float(conf or 0.5)
            if score > 0.01:
                scored.append((score, source_ref, chunk_text))

        scored.sort(key=lambda x: x[0], reverse=True)
        return [
            {"score": s, "source": src, "text": txt}
            for s, src, txt in scored[:top_k]
        ]

    def _analyze_images(self, images: List[str], user_text: str) -> Tuple[str, str]:
        valid = []
        for img in images[:3]:
            if not isinstance(img, str):
                continue
            b64 = self._image_ref_to_base64(img)
            if b64:
                valid.append(b64)

        if not valid:
            return "", "No valid image payload found"

        prompt = (
            "Analyze this image carefully for visible UI states, errors, key text, and actionable facts. "
            "Respond with short bullet points. User question: " + (user_text or "")
        )
        payload = {
            "model": self.valves.VISION_MODEL,
            "messages": [{"role": "user", "content": prompt, "images": valid}],
            "stream": False,
        }
        try:
            data = self._http_json(
                f"{self.valves.OLLAMA_BASE_URL}/api/chat", payload, timeout=120
            )
            msg = (data.get("message") or {}).get("content", "").strip()
            return msg, ""
        except Exception as exc:
            return "", f"Vision analysis failed: {exc}"

    def _compose_system_context(
        self,
        mode: str,
        retrieval_context: str,
        vision_context: str,
        file_refs: List[str],
        parse_notes: List[str],
        uncertainty_notes: List[str],
    ) -> str:
        sections = [
            "You are Getouch Smart Assistant running in a routed orchestration mode.",
            f"Routing mode: {mode}",
            "Response format rules:",
            "1) Start with a short section titled 'Ringkasan' (2-5 concise bullet points).",
            "2) Then provide a detailed section titled 'Perincian'.",
            "3) If documents are used, include supporting references in 'Sources:' section.",
            "4) Include uncertainty notes if extraction confidence is low.",
            "5) Do not fabricate file content.",
        ]

        if file_refs:
            sections.append("Files used: " + ", ".join(file_refs))
        if retrieval_context:
            sections.append("Retrieved context:\n" + retrieval_context)
        if vision_context:
            sections.append("Vision analysis notes:\n" + vision_context)
        if parse_notes:
            sections.append("Parser notes:\n- " + "\n- ".join(parse_notes))
        if uncertainty_notes:
            sections.append("Uncertainty:\n- " + "\n- ".join(uncertainty_notes))
        return "\n\n".join(sections)

    def _last_user_text(self, messages: List[Dict[str, Any]]) -> str:
        for m in reversed(messages):
            if m.get("role") != "user":
                continue
            content = m.get("content")
            if isinstance(content, str):
                return content
            if isinstance(content, list):
                parts = []
                for p in content:
                    if isinstance(p, dict) and p.get("type") in {"text", "input_text"}:
                        parts.append(str(p.get("text") or p.get("content") or ""))
                return "\n".join(parts)
        return ""

    def _http_json(self, url: str, payload: Dict[str, Any], timeout: int = 30) -> Dict[str, Any]:
        raw = json.dumps(payload).encode("utf-8")
        req = request.Request(
            url,
            data=raw,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        try:
            with request.urlopen(req, timeout=timeout) as resp:
                body = resp.read().decode("utf-8")
        except HTTPError as exc:
            err_body = ""
            try:
                err_body = exc.read().decode("utf-8", errors="replace")
            except Exception:
                err_body = ""
            detail = err_body[:500].strip()
            if detail:
                raise RuntimeError(f"HTTP {exc.code} from {url}: {detail}")
            raise RuntimeError(f"HTTP {exc.code} from {url}")
        return json.loads(body)

    def _complete_ollama_chat(
        self,
        target_model: str,
        messages: List[Dict[str, Any]],
        options: Dict[str, Any],
        timeout: int = 240,
    ) -> str:
        payload = {
            "model": target_model,
            "messages": messages,
            "stream": False,
        }
        if options:
            payload["options"] = options
        data = self._http_json(
            f"{self.valves.OLLAMA_BASE_URL}/api/chat",
            payload,
            timeout=timeout,
        )
        return (data.get("message") or {}).get("content", "")

    def _stream_ollama_chat(
        self,
        target_model: str,
        messages: List[Dict[str, Any]],
        options: Dict[str, Any],
        web_image_urls: List[str],
        place_cards: Optional[List[Dict[str, str]]] = None,
        sources: Optional[List[Dict[str, str]]] = None,
    ) -> Iterator[str]:
        payload = {
            "model": target_model,
            "messages": messages,
            "stream": True,
        }
        if options:
            payload["options"] = options

        raw = json.dumps(payload).encode("utf-8")
        req = request.Request(
            f"{self.valves.OLLAMA_BASE_URL}/api/chat",
            data=raw,
            headers={"Content-Type": "application/json"},
            method="POST",
        )

        if place_cards:
            yield self._format_place_cards(place_cards) + "\n\n"
        elif web_image_urls:
            yield self._format_image_markdown(web_image_urls) + "\n\n"

        if sources:
            yield f"\n> [\U0001f50d Read {len(sources)} web sources](#getouch-sources)\n\n"

        with request.urlopen(req, timeout=300) as resp:
            for line in resp:
                if not line:
                    continue
                try:
                    obj = json.loads(line.decode("utf-8", errors="ignore").strip())
                except Exception:
                    continue
                msg = (obj.get("message") or {}).get("content", "")
                if msg:
                    yield msg

        if sources:
            safe = json.dumps(sources, ensure_ascii=False)
            yield f"\n\n```getouch-sources\n{safe}\n```\n"

    def _to_ollama_messages(self, messages: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        out: List[Dict[str, Any]] = []
        for m in messages:
            role = str(m.get("role") or "user")
            content = m.get("content")
            text_parts: List[str] = []
            images: List[str] = []

            if isinstance(content, str):
                text_parts.append(content)
            elif isinstance(content, list):
                for part in content:
                    if not isinstance(part, dict):
                        continue
                    ptype = str(part.get("type") or "").lower()
                    if ptype in {"text", "input_text"}:
                        t = part.get("text") or part.get("content")
                        if t:
                            text_parts.append(str(t))
                    elif ptype in {"image_url", "input_image", "image"}:
                        img_ref = ""
                        if isinstance(part.get("image_url"), dict):
                            img_ref = str(part.get("image_url", {}).get("url") or "")
                        elif part.get("image"):
                            img_ref = str(part.get("image") or "")
                        if img_ref:
                            b64 = self._image_ref_to_base64(img_ref)
                            if b64:
                                images.append(b64)

            for img in self._coerce_list(m.get("images")):
                if isinstance(img, str):
                    b64 = self._image_ref_to_base64(img)
                    if b64:
                        images.append(b64)

            msg_obj: Dict[str, Any] = {
                "role": role,
                "content": "\n".join([p for p in text_parts if p]).strip() or " ",
            }
            if images and role == "user":
                # Ollama expects base64 image data without data-uri prefix.
                msg_obj["images"] = images[: max(1, int(self.valves.MAX_IMAGES_PER_TURN))]
            out.append(msg_obj)
        return out

    def _image_ref_to_base64(self, ref: str) -> str:
        if not ref:
            return ""
        if ref.startswith("data:image") and ";base64," in ref:
            return ref.split(";base64,", 1)[1]
        if ref.startswith("/"):
            ref = f"{self.valves.OPENWEBUI_BASE_URL.rstrip('/')}{ref}"
        if ref.startswith("http://") or ref.startswith("https://"):
            try:
                req = request.Request(ref, headers={"User-Agent": "Mozilla/5.0"})
                with request.urlopen(req, timeout=20) as resp:
                    content_type = str(resp.headers.get("Content-Type") or "").lower()
                    img_bytes = resp.read()
                if content_type and not content_type.startswith("image/"):
                    return ""
                if imghdr.what(None, img_bytes) is None:
                    return ""
                return base64.b64encode(img_bytes).decode("utf-8")
            except Exception:
                return ""
        # Not a supported image reference payload for direct model usage.
        return ""

    def _should_search_place_images(self, user_text: str, mode: str) -> bool:
        """Detect queries that benefit from per-place image cards (travel, recommendations, itineraries)."""
        if mode not in {"text_only", "tool_required"}:
            return False
        q = (user_text or "").lower()
        return bool(
            re.search(
                r"\b(i[lt]*[ie]n?[ae]?rar[iy]|i[lt]+ernary|travel|trip|tempat|lawat|melawat|jalan[- ]?jalan"
                r"|guide|plan|visit|sightseeing|holiday|vacation|percutian|cuti"
                r"|recommend|cadang|saran|suggest|best\s+place|top\s+\d+"
                r"|things?\s+to\s+do|what\s+to\s+see|where\s+to\s+go"
                r"|makan\s+sedap|restoran|restaurant|cafe|hotel|resort"
                r"|destinasi|pelancongan|backpack|roadtrip|staycation"
                r"|\d+\s*(?:hari|malam|night|day))",
                q,
            )
        )

    def _should_search_web_images(self, user_text: str, mode: str) -> bool:
        if mode not in {"text_only", "tool_required"}:
            return False
        q = (user_text or "").lower()
        return bool(
            re.search(
                r"\b(itinerary|travel|trip|tempat|guide|plan|visit|sightseeing)\b",
                q,
            )
        )

    # ── Travel itinerary planner ─────────────────────────────

    def _is_travel_itinerary(self, user_text: str) -> bool:
        """Detect queries that specifically request a travel itinerary or trip plan."""
        q = (user_text or "").lower()
        if re.search(
            r"\b(i[lt]*[ie]n?[ae]?rar[iy]|i[lt]+ernary|rencana\s+perjalanan"
            r"|perancangan\s+(?:perjalanan|trip|cuti|percutian))\b",
            q,
        ):
            return True
        if re.search(r"\b(?:buat|buatkan|create|make|plan|tolong)\b", q) and re.search(
            r"\b(?:perjalanan|trip|travel|holiday|vacation|percutian|cuti)\b", q
        ):
            return True
        if re.search(r"\d+\s*(?:hari|days?|malam|nights?|mlm)\b", q) and re.search(
            r"\b(?:ke|to|in|at)\s+\w", q
        ):
            return True
        return False

    def _parse_trip_request(self, user_text: str) -> Dict[str, Any]:
        """Extract trip duration from user query."""
        q = (user_text or "").lower()
        m = re.search(r"(\d+)\s*(?:hari|days?|malam|nights?|mlm)", q)
        days = int(m.group(1)) if m else 3
        is_nights = bool(m and re.search(r"malam|nights?|mlm", m.group(0))) if m else False
        if is_nights:
            days += 1
        days = max(1, min(days, 14))
        return {"days": days, "nights": max(1, days - 1)}

    def _search_web_sources(
        self, destination: str, trip_info: Dict[str, Any]
    ) -> List[Dict[str, str]]:
        """Search SearXNG for text sources about a travel destination."""
        days = trip_info.get("days", 3)
        queries = [
            f"{destination} travel itinerary {days} days guide",
            f"{destination} attractions food transport tips",
        ]
        sources: List[Dict[str, str]] = []
        seen_domains: set = set()
        for query in queries:
            try:
                q = quote_plus(query[:200])
                url = (
                    f"{self.valves.SEARXNG_BASE_URL}/search?format=json"
                    f"&categories=general&language=en&q={q}"
                )
                req_obj = request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
                raw = request.urlopen(req_obj, timeout=15).read().decode(
                    "utf-8", errors="ignore"
                )
                data = json.loads(raw)
                for r in (data.get("results") or [])[:10]:
                    if not isinstance(r, dict):
                        continue
                    page_url = r.get("url", "")
                    if not page_url:
                        continue
                    domain = self._extract_domain(page_url)
                    if domain in seen_domains:
                        continue
                    seen_domains.add(domain)
                    snippet = (r.get("content") or "")[:300].strip()
                    if not snippet:
                        continue
                    sources.append({
                        "id": f"src_{len(sources) + 1}",
                        "title": (r.get("title") or domain)[:120],
                        "url": page_url,
                        "domain": domain,
                        "snippet": snippet,
                    })
                    if len(sources) >= 8:
                        break
            except Exception as exc:
                self._log("web_source_error", {"query": query[:80], "error": str(exc)})
            if len(sources) >= 8:
                break
        self._log("web_sources", {"destination": destination, "found": len(sources)})
        return sources

    @staticmethod
    def _extract_domain(url: str) -> str:
        try:
            host = urlparse(url).netloc
            return host[4:] if host.startswith("www.") else host
        except Exception:
            return ""

    def _compose_travel_planner_prompt(
        self,
        user_message: str,
        trip_info: Dict[str, Any],
        pois: List[Dict[str, str]],
        web_sources: List[Dict[str, str]],
    ) -> str:
        """Build a comprehensive travel planner system prompt."""
        days = trip_info["days"]
        nights = trip_info["nights"]
        destination = ""
        if pois:
            city = pois[0].get("city", "")
            country = pois[0].get("country", "")
            destination = f"{city}, {country}".strip(", ")
        poi_names = ", ".join(p["display_name"] for p in pois) if pois else "use your knowledge"

        source_ctx = ""
        if web_sources:
            parts = []
            for i, s in enumerate(web_sources, 1):
                parts.append(f"[{i}] {s['title']} ({s['domain']})\n{s['snippet']}")
            source_ctx = (
                "\nWEB RESEARCH CONTEXT (cite with [N] inline where relevant):\n"
                + "\n\n".join(parts)
                + "\n"
            )

        return (
            "You are an expert travel planner AI. Produce a comprehensive, well-researched "
            "travel itinerary that feels like a professional travel guide.\n\n"
            f"DESTINATION: {destination or 'as stated by user'}\n"
            f"TRIP LENGTH: {days} days / {nights} nights\n"
            f"KNOWN ATTRACTIONS: {poi_names}\n"
            f"{source_ctx}\n"
            "RESPONSE LANGUAGE: Match the user's language. If the user wrote in Malay, "
            "respond primarily in Malay/Bahasa.\n\n"
            "YOUR ANSWER MUST INCLUDE ALL SECTIONS BELOW IN THIS EXACT ORDER:\n\n"
            f"## \U0001f5fa\ufe0f Rencana Perjalanan {destination} {days} Hari\n\n"
            "Write 2-3 engaging intro sentences about the destination and trip.\n\n"
            "### \U0001f4cb Ringkasan Perjalanan\n\n"
            "| Hari | Tema | Aktiviti Utama |\n"
            "|------|------|----------------|\n"
            f"(one row per day for all {days} days)\n\n"
            "### \U0001f4c5 Hari 1: [Theme]\n\n"
            "**\U0001f305 Pagi**\n"
            "- Specific activity with place name and practical detail\n\n"
            "**\U0001f324\ufe0f Tengah Hari / Petang**\n"
            "- Afternoon activities\n\n"
            "**\U0001f319 Malam**\n"
            "- Evening plan with dinner suggestion\n\n"
            "> \U0001f4a1 **Tip:** Practical tip for this day\n\n"
            f"(Repeat for ALL {days} days. Day 1 = arrival, Day {days} = departure.)\n\n"
            "### \U0001f35c Cadangan Makanan\n\n"
            "5-8 must-try local foods:\n"
            "- **Food name** \u2014 description, where to try (area/street)\n\n"
            "### \U0001f697 Tips Pengangkutan\n\n"
            "How to reach destination, local transport, parking, walkability.\n\n"
            "### \U0001f3e8 Kawasan Penginapan\n\n"
            "2-3 areas:\n"
            "- **Area** \u2014 why, proximity, budget level\n\n"
            "### \U0001f4b0 Anggaran Perbelanjaan\n\n"
            "| Item | Anggaran (per orang/hari) |\n"
            "|------|---------------------------|\n"
            "| Penginapan | ... |\n"
            "| Makanan | ... |\n"
            "| Pengangkutan | ... |\n"
            "| Tiket/Aktiviti | ... |\n"
            "| **Jumlah** | **...** |\n\n"
            "### \U0001f4dd Tips Praktikal\n\n"
            "5-8 tips: weather, timing, clothing, booking, cultural etiquette, safety.\n\n"
            "CRITICAL RULES:\n"
            "- Group attractions by area to minimize backtracking\n"
            f"- Day 1 lighter (arrival), Day {days} lighter (departure)\n"
            "- Mix: heritage, food, scenic, shopping, local exploration\n"
            "- Be SPECIFIC: real place names, street names, local food names\n"
            "- Do NOT invent fake museums, districts, or transport systems\n"
            "- When uncertain, use softer wording\n"
            "- Avoid generic filler \u2014 give actionable details\n"
            "- Do NOT include a Sources/Sumber section \u2014 added automatically\n"
            "- Cite web research inline as [1], [2] where applicable"
        )

    def _handle_travel_itinerary(
        self,
        user_message: str,
        target_model: str,
        ollama_messages: List[Dict[str, Any]],
        body: dict,
        started: float,
    ) -> Union[str, Iterator[str]]:
        """Enhanced travel itinerary flow with web enrichment and structured output."""
        trip_info = self._parse_trip_request(user_message)

        pois = self._extract_place_names(user_message)
        destination = ""
        if pois:
            city = pois[0].get("city", "")
            country = pois[0].get("country", "")
            destination = f"{city}, {country}".strip(", ")

        web_sources = self._search_web_sources(
            destination or user_message[:60], trip_info
        )

        place_cards: List[Dict[str, str]] = []
        for poi in pois:
            result = self._search_place_image(poi)
            if result:
                place_cards.append({
                    "name": poi["display_name"],
                    "image_url": result["image_url"],
                })

        travel_prompt = self._compose_travel_planner_prompt(
            user_message=user_message,
            trip_info=trip_info,
            pois=pois,
            web_sources=web_sources,
        )

        if ollama_messages and ollama_messages[0].get("role") == "system":
            ollama_messages[0]["content"] = travel_prompt
        else:
            ollama_messages.insert(0, {"role": "system", "content": travel_prompt})

        options = {"temperature": 0.7, "num_predict": 8192}
        stream_requested = bool(body.get("stream", False))

        self._log("travel_planner", {
            "destination": destination,
            "days": trip_info["days"],
            "pois": len(pois),
            "sources": len(web_sources),
            "cards": len(place_cards),
            "latency_ms": int((time.time() - started) * 1000),
        })

        if stream_requested:
            return self._stream_ollama_chat(
                target_model=target_model,
                messages=ollama_messages,
                options=options,
                web_image_urls=[],
                place_cards=place_cards,
                sources=web_sources,
            )

        try:
            content = self._complete_ollama_chat(
                target_model=target_model,
                messages=ollama_messages,
                options=options,
                timeout=300,
            )
            if not isinstance(content, str) or not content.strip():
                content = "Unable to generate travel itinerary. Please try again."
            parts = []
            if place_cards:
                parts.append(self._format_place_cards(place_cards))
            if web_sources:
                parts.append(f"> [\U0001f50d Read {len(web_sources)} web sources](#getouch-sources)")
            parts.append(content)
            if web_sources:
                safe = json.dumps(web_sources, ensure_ascii=False)
                parts.append(f"```getouch-sources\n{safe}\n```")
            return "\n\n".join(parts)
        except Exception as exc:
            self._log("travel_planner_error", {"error": str(exc)})
            return "Travel planner encountered an error. Please try again."

    def _format_sources_section(self, sources: List[Dict[str, str]]) -> str:
        """Render sources as a markdown footer with clickable links."""
        if not sources:
            return ""
        lines = ["---", "", "### \U0001f4da Sumber", ""]
        for i, s in enumerate(sources, 1):
            title = s.get("title", s.get("domain", "Source"))
            url = s.get("url", "")
            domain = s.get("domain", "")
            if url:
                lines.append(f"{i}. [{title}]({url}) \u2014 *{domain}*")
            else:
                lines.append(f"{i}. {title}")
        return "\n".join(lines)

    def _extract_place_names(self, user_text: str) -> List[Dict[str, str]]:
        """Quick LLM call to extract destination + POIs with canonical names.

        Returns list of dicts: {display_name, canonical_name, city, country}
        """
        prompt = (
            "You are a travel data extractor. From the user query, extract:\n"
            "1. The destination city and country\n"
            "2. 4-6 famous tourist attractions AT that destination\n\n"
            "Output EXACTLY in this format, one per line:\n"
            "DESTINATION: <City>, <Country>\n"
            "POI: <Display Name> | <Full Canonical Name including city and country>\n\n"
            "Rules:\n"
            "- Normalize misspellings (malaka→Melaka, kunning→Kunming)\n"
            "- Canonical name MUST include city and country\n"
            "- Include landmark type when useful (fort, temple, lake, church)\n\n"
            "Example for 'buat itinerary ke malaka 3 hari':\n"
            "DESTINATION: Melaka, Malaysia\n"
            "POI: A Famosa | A Famosa Fort, Melaka, Malaysia\n"
            "POI: Christ Church | Christ Church, Melaka, Malaysia\n"
            "POI: Jonker Walk | Jonker Street, Melaka, Malaysia\n"
            "POI: Stadthuys | The Stadthuys, Melaka, Malaysia\n\n"
            f"User query: {user_text[:300]}"
        )
        try:
            content = self._complete_ollama_chat(
                target_model=self.valves.TEXT_MODEL,
                messages=[{"role": "user", "content": prompt}],
                options={"temperature": 0.2, "num_predict": 250},
                timeout=25,
            )
            destination = ""
            pois: List[Dict[str, str]] = []
            for line in content.strip().split("\n"):
                line = line.strip()
                if line.upper().startswith("DESTINATION:"):
                    dest_raw = line.split(":", 1)[1].strip()
                    parts = [p.strip() for p in dest_raw.split(",")]
                    destination = dest_raw
                    city = parts[0] if parts else ""
                    country = parts[1] if len(parts) > 1 else ""
                elif line.upper().startswith("POI:"):
                    poi_raw = line.split(":", 1)[1].strip()
                    if "|" in poi_raw:
                        display, canonical = [p.strip() for p in poi_raw.split("|", 1)]
                    else:
                        display = poi_raw.strip().strip("0123456789.-) *")
                        canonical = f"{display}, {destination}" if destination else display
                    if display and len(display) > 2:
                        # Ensure canonical always has destination context
                        canonical_lower = canonical.lower()
                        if destination:
                            dest_parts = [p.strip().lower() for p in destination.split(",")]
                            if not any(dp in canonical_lower for dp in dest_parts if dp):
                                canonical = f"{canonical}, {destination}"
                        pois.append({
                            "display_name": display,
                            "canonical_name": canonical,
                            "city": city if 'city' in dir() else "",
                            "country": country if 'country' in dir() else "",
                        })
            self._log("place_extract", {
                "destination": destination,
                "pois": [p["display_name"] for p in pois],
                "canonical": [p["canonical_name"] for p in pois],
            })
            return pois[:int(self.valves.PLACE_IMAGE_LIMIT)]
        except Exception as exc:
            self._log("place_extract_error", {"error": str(exc)})
            return []

    def _score_image_candidate(
        self, poi: Dict[str, str], candidate: Dict[str, Any]
    ) -> float:
        """Score an image candidate based on metadata relevance to the POI.

        Returns a score from -1.0 to 1.0.  Higher = more relevant.
        """
        title = str(candidate.get("title") or "").lower()
        snippet = str(candidate.get("content") or "").lower()
        src_url = str(candidate.get("img_src") or candidate.get("url") or "").lower()
        meta = f"{title} {snippet} {src_url}"

        canonical = poi.get("canonical_name", "").lower()
        display = poi.get("display_name", "").lower()
        city = poi.get("city", "").lower()
        country = poi.get("country", "").lower()

        score = 0.0

        # Positive signals: destination context in metadata
        if city and city in meta:
            score += 0.35
        if country and country in meta:
            score += 0.20
        # Match display or canonical name tokens
        display_tokens = [t for t in display.split() if len(t) > 2]
        match_count = sum(1 for t in display_tokens if t in meta)
        if display_tokens:
            score += 0.30 * (match_count / len(display_tokens))

        # Negative signals: obviously wrong landmarks
        wrong_landmarks = [
            "taj mahal", "eiffel tower", "colosseum", "statue of liberty",
            "big ben", "sydney opera", "great wall", "machu picchu",
            "pyramids", "petra jordan", "angkor wat", "christ the redeemer",
        ]
        for wl in wrong_landmarks:
            if wl in meta and wl not in canonical:
                score -= 0.80

        # Negative: wrong country in metadata
        wrong_countries = [
            "india", "france", "italy", "usa", "united states", "england",
            "australia", "peru", "egypt", "jordan", "brazil", "cambodia",
        ]
        if country:
            for wc in wrong_countries:
                if wc in meta and wc != country:
                    score -= 0.40

        return max(-1.0, min(1.0, score))

    def _search_place_image(self, poi: Dict[str, str]) -> Optional[Dict[str, Any]]:
        """Search SearXNG images for a POI using its canonical name.

        Returns {image_url, score, source_title} or None.
        """
        canonical = poi.get("canonical_name", poi.get("display_name", ""))
        try:
            q = quote_plus(f"{canonical} landmark photo")
            url = (
                f"{self.valves.SEARXNG_BASE_URL}/search?format=json&categories=images"
                f"&language=en&q={q}"
            )
            req = request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            raw = request.urlopen(req, timeout=15).read().decode("utf-8", errors="ignore")
            data = json.loads(raw)
            results = data.get("results") or []

            best_url: Optional[str] = None
            best_score = -1.0
            best_title = ""
            candidates_checked = 0

            for r in results[:12]:
                if not isinstance(r, dict):
                    continue
                img_url = r.get("img_src") or r.get("thumbnail")
                if not isinstance(img_url, str) or not img_url.startswith("http"):
                    continue
                candidates_checked += 1
                sc = self._score_image_candidate(poi, r)
                if sc > best_score:
                    best_score = sc
                    best_url = img_url
                    best_title = str(r.get("title") or "")

            min_threshold = 0.15
            self._log("image_search", {
                "poi": poi.get("display_name"),
                "canonical": canonical,
                "query": f"{canonical} landmark photo",
                "candidates_checked": candidates_checked,
                "best_score": round(best_score, 3),
                "best_title": best_title[:80],
                "passed": best_score >= min_threshold,
            })

            if best_url and best_score >= min_threshold:
                return {"image_url": best_url, "score": best_score, "source_title": best_title}
        except Exception as exc:
            self._log("image_search_error", {"poi": poi.get("display_name"), "error": str(exc)})
        return None

    def _build_place_cards(self, user_text: str) -> List[Dict[str, str]]:
        """Extract places and find a validated image for each."""
        pois = self._extract_place_names(user_text)
        if not pois:
            return []

        cards: List[Dict[str, str]] = []
        for poi in pois:
            result = self._search_place_image(poi)
            if result:
                cards.append({
                    "name": poi["display_name"],
                    "image_url": result["image_url"],
                })
            # Skip POIs with no confident image match

        self._log("place_cards", {"requested": len(pois), "found": len(cards)})
        return cards

    def _format_place_cards(self, cards: List[Dict[str, str]]) -> str:
        """Render place cards as a markdown image table."""
        if not cards:
            return ""
        img_cells = []
        name_cells = []
        for c in cards:
            name = c["name"].replace("|", "\\|")
            url = c["image_url"].replace("(", "%28").replace(")", "%29")
            img_cells.append(f"![{name}]({url})")
            name_cells.append(f"**{name}**")
        img_row = "| " + " | ".join(img_cells) + " |"
        sep_row = "| " + " | ".join(":---:" for _ in cards) + " |"
        name_row = "| " + " | ".join(name_cells) + " |"
        return f"{img_row}\n{sep_row}\n{name_row}"

    def _search_web_images(self, query: str, limit: int) -> List[str]:
        try:
            q = quote_plus(query[:180])
            url = (
                f"{self.valves.SEARXNG_BASE_URL}/search?format=json&categories=images"
                f"&language=ms&q={q}"
            )
            req = request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            raw = request.urlopen(req, timeout=20).read().decode("utf-8", errors="ignore")
            data = json.loads(raw)
            results = data.get("results") or []
            urls: List[str] = []
            for r in results:
                if not isinstance(r, dict):
                    continue
                cand = r.get("img_src") or r.get("thumbnail") or r.get("url")
                if isinstance(cand, str) and cand.startswith("http"):
                    urls.append(cand)
                if len(urls) >= max(1, limit):
                    break
            return urls
        except Exception:
            return []

    def _format_image_markdown(self, urls: List[str]) -> str:
        if not urls:
            return ""
        lines = ["Visual references:"]
        for idx, u in enumerate(urls, start=1):
            lines.append(f"![reference-{idx}]({u})")
        return "\n".join(lines)

    def _keyword_overlap(self, query: str, text: str) -> float:
        q_tokens = {t for t in re.findall(r"[a-z0-9]+", query.lower()) if len(t) > 2}
        t_tokens = {t for t in re.findall(r"[a-z0-9]+", text.lower()) if len(t) > 2}
        if not q_tokens or not t_tokens:
            return 0.0
        inter = len(q_tokens.intersection(t_tokens))
        return inter / max(len(q_tokens), 1)

    def _cosine(self, a: List[float], b: List[float]) -> float:
        if not a or not b or len(a) != len(b):
            return 0.0
        dot = sum(x * y for x, y in zip(a, b))
        na = math.sqrt(sum(x * x for x in a))
        nb = math.sqrt(sum(y * y for y in b))
        if na == 0 or nb == 0:
            return 0.0
        return dot / (na * nb)

    def _coerce_list(self, v: Any) -> List[Any]:
        if v is None:
            return []
        if isinstance(v, list):
            return v
        return [v]

    def _db_execute(self, sql: str, params: Tuple[Any, ...]):
        with self.lock:
            con = sqlite3.connect(self.db_path)
            cur = con.cursor()
            cur.execute(sql, params)
            con.commit()
            con.close()

    def _db_fetchone(self, sql: str, params: Tuple[Any, ...]):
        with self.lock:
            con = sqlite3.connect(self.db_path)
            cur = con.cursor()
            cur.execute(sql, params)
            row = cur.fetchone()
            con.close()
            return row

    def _db_fetchall(self, sql: str, params: Tuple[Any, ...]):
        with self.lock:
            con = sqlite3.connect(self.db_path)
            cur = con.cursor()
            cur.execute(sql, params)
            rows = cur.fetchall()
            con.close()
            return rows

    def _log(self, event: str, payload: Dict[str, Any]):
        if not self.valves.ENABLE_STRUCTURED_LOGS:
            return
        line = {
            "ts": int(time.time()),
            "event": event,
            "payload": payload,
        }
        try:
            print(json.dumps(line, ensure_ascii=True))
        except Exception:
            print(f"{event}: {payload}")
            print(traceback.format_exc())
